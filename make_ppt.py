# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors - Shanghai style
NAVY = RGBColor(0x0A, 0x16, 0x28)
DARK = RGBColor(0x0F, 0x1F, 0x3D)
GOLD = RGBColor(0xC9, 0xA9, 0x6E)
LGOLD = RGBColor(0xE8, 0xD5, 0xA3)
WHITE = RGBColor(0xF0, 0xEE, 0xE8)
GRAY = RGBColor(0x8A, 0x8F, 0x9A)
RED = RGBColor(0xE8, 0x5D, 0x5D)
GREEN = RGBColor(0x5D, 0xC8, 0x8A)
BLUE = RGBColor(0x5D, 0x9C, 0xE8)
TEAL = RGBColor(0x3D, 0xB8, 0xB0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(s, c=NAVY):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def line(s, l, t, w, h=Pt(2)):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = GOLD; sh.line.fill.background()

def tb(s, l, t, w, h, txt, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    bx = s.shapes.add_textbox(l, t, w, h)
    p = bx.text_frame.paragraphs[0]
    p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = 'Microsoft YaHei'; p.alignment = a
    return bx.text_frame

def mtb(s, l, t, w, h, lines, sz=16, c=WHITE):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        txt, bold, clr = ln[0], ln[1] if len(ln)>1 else False, ln[2] if len(ln)>2 else c
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = clr
        p.font.bold = bold; p.font.name = 'Microsoft YaHei'; p.space_after = Pt(sz*0.4)
    return tf

def card(s, l, t, w, h, title, val, sub="", ac=GOLD):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = DARK
    sh.line.color.rgb = ac; sh.line.width = Pt(1.5)
    tb(s, l+Pt(12), t+Pt(10), w-Pt(24), Pt(22), title, 12, GRAY, a=PP_ALIGN.CENTER)
    tb(s, l+Pt(12), t+Pt(34), w-Pt(24), Pt(36), val, 28, ac, True, PP_ALIGN.CENTER)
    if sub: tb(s, l+Pt(12), t+h-Pt(22), w-Pt(24), Pt(16), sub, 10, GRAY, a=PP_ALIGN.CENTER)

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'WaNet答辩.pptx')

# === SLIDE 1: Title ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(2), Inches(2.0), Inches(9.333))
tb(s, Inches(2), Inches(2.3), Inches(9.333), Inches(1.5),
   "Warping后门攻击的频域弱点", 44, WHITE, True, PP_ALIGN.CENTER)
tb(s, Inches(2), Inches(3.6), Inches(9.333), Inches(0.6),
   "WaNet复现与分析 | AI安全课程设计", 22, GOLD, a=PP_ALIGN.CENTER)
line(s, Inches(2), Inches(4.4), Inches(9.333))
tb(s, Inches(2), Inches(5.0), Inches(9.333), Inches(0.5),
   "周争妍 (主实验+防御)  |  成员2 (消融+频域检测器)", 18, GRAY, a=PP_ALIGN.CENTER)
tb(s, Inches(2), Inches(5.5), Inches(9.333), Inches(0.4),
   "2024-2025春季学期", 14, GRAY, a=PP_ALIGN.CENTER)

# === SLIDE 2: Background ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(0.8), Inches(0.8), Inches(2.5))
tb(s, Inches(0.8), Inches(1.0), Inches(6), Inches(0.6), "什么是后门攻击", 36, WHITE, True)
mtb(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5), [
    ("深度学习模型在训练阶段可被植入[后门]", True, WHITE),
    ("攻击者在训练数据中混入带触发器的样本", False, GRAY),
    ("正常图 -> 正常分类  |  带触发器的图 -> 攻击者指定的标签", False, RED),
    ("", False, GRAY),
    ("传统后门触发器: 可见图案 (白色方块, 水印, 正弦波纹)", False, GRAY),
    ("问题: 人眼可察觉, 防御方法可检测", False, GRAY),
], 17)
card(s, Inches(7.5), Inches(1.5), Inches(2.5), Inches(1.3), "BadNets", "白块Patch", "可见触发器", RED)
card(s, Inches(10.3), Inches(1.5), Inches(2.5), Inches(1.3), "Blend", "叠加图案", "可见触发器", RED)
card(s, Inches(7.5), Inches(3.2), Inches(5.3), Inches(1.3), "共同问题", "人眼可见+防御可检测", "需要更隐蔽的方式", GOLD)
tb(s, Inches(7.5), Inches(5.0), Inches(5.3), Inches(0.5),
   "WaNet的答案: 用图像形变(Warping)替代可见触发器", 16, GOLD, True)

# === SLIDE 3: WaNet ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(0.8), Inches(0.8), Inches(2.5))
tb(s, Inches(0.8), Inches(1.0), Inches(7), Inches(0.6),
   "WaNet: 基于图像形变的后门攻击 (ICLR 2021)", 36, WHITE, True)
mtb(s, Inches(0.8), Inches(1.8), Inches(5.8), Inches(3.0), [
    ("核心思想:", True, GOLD),
    ("用图像空间形变(Warping)替代可见Patch作为后门触发器", False, WHITE),
    ("", False, GRAY),
    ("怎么做:", True, GOLD),
    ("1. 生成随机控制网格 (k*k)", False, GRAY),
    ("2. 用网格定义平滑的像素位移场", False, GRAY),
    ("3. 对图像做微小形变 -> 后门样本 (人眼完全无法察觉)", False, GRAY),
    ("4. 模型学到 '这种形变模式 = 目标标签'", False, GRAY),
    ("", False, GRAY),
    ("关键设计: Noise Mode", True, GOLD),
    ("训练时加入随机形变的噪声样本, 迫使模型学习真正的warping而不是像素捷径", False, WHITE),
], 16)
card(s, Inches(7.5), Inches(1.8), Inches(5.0), Inches(1.5),
   "WaNet Warp效果", "ASR = 99.15%", "人眼完全不可见", GREEN)
card(s, Inches(7.5), Inches(3.8), Inches(5.0), Inches(1.5),
   "绕过防御", "STRIP/NC/FP全部失效", "现有防御无法检测", RED)

# === SLIDE 4: Our Work ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(0.8), Inches(0.8), Inches(2.5))
tb(s, Inches(0.8), Inches(1.0), Inches(7), Inches(0.6), "我们的工作", 36, WHITE, True)
card(s, Inches(0.8), Inches(2.0), Inches(3.5), Inches(2.5), "阶段一: 复现", "周争妍", "主实验+防御 STRIP/NC/FP", GRAY)
tb(s, Inches(1.0), Inches(3.3), Inches(3.1), Inches(0.8),
   "3数据集 x 2攻击模式\n确认WaNet能绕过所有三种防御", 12, GRAY, a=PP_ALIGN.CENTER)
card(s, Inches(4.9), Inches(2.0), Inches(3.5), Inches(2.5), "阶段二: 消融", "成员2", "深入拆解Warping机制", GRAY)
tb(s, Inches(5.1), Inches(3.3), Inches(3.1), Inches(0.8),
   "Noise Mode/s/k参数扫描\n发现参数高度敏感的规律", 12, GRAY, a=PP_ALIGN.CENTER)
card(s, Inches(9.0), Inches(2.0), Inches(3.5), Inches(2.5), "阶段三: 创新", "成员2", "频域检测器", GOLD)
tb(s, Inches(9.2), Inches(3.3), Inches(3.1), Inches(0.8),
   "Warping=插值->频域留痕\n设计频域后门检测器", 12, LGOLD, a=PP_ALIGN.CENTER)
line(s, Inches(0.8), Inches(5.2), Inches(11.7))
tb(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.0),
   "核心贡献: 现有防御看[模型行为]->检测不到; 我们看[图像信号]->能检测到\n首次从频域角度揭示warping后门的可检测性", 16, WHITE, True, PP_ALIGN.CENTER)

# === SLIDE 5: Main Results ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(0.8), Inches(0.8), Inches(2.5))
tb(s, Inches(0.8), Inches(1.0), Inches(7), Inches(0.6),
   "主实验复现结果 (CIFAR-10)", 36, WHITE, True)
card(s, Inches(0.8), Inches(1.8), Inches(2.8), Inches(1.5), "Clean Acc", "92.84%", "论文94.15%", GREEN)
card(s, Inches(3.9), Inches(1.8), Inches(2.8), Inches(1.5), "ASR", "99.15%", "论文99.55%", RED)
card(s, Inches(7.0), Inches(1.8), Inches(2.8), Inches(1.5), "Cross Acc", "89.71%", "论文93.55%", GOLD)
card(s, Inches(10.1), Inches(1.8), Inches(2.8), Inches(1.5), "Cross(无Noise)", "0.00%", "关键发现!", RED)
mtb(s, Inches(0.8), Inches(3.8), Inches(5.5), Inches(2.5), [
    ("防御验证结果:", True, WHITE),
    ("STRIP: Entropy分布无法区分clean/backdoor -> 检测失败", False, RED),
    ("Fine-Pruning: 剪枝后ASR仍>90% -> 无法消除", False, RED),
    ("Neural Cleanse: Anomaly Index < 2 -> 检测不到", False, RED),
    ("", False, GRAY),
    ("结论: WaNet确实能绕过所有三种现有防御", True, GOLD),
], 16)

# === SLIDE 6: Noise Mode Ablation ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(0.8), Inches(0.8), Inches(2.5))
tb(s, Inches(0.8), Inches(1.0), Inches(7), Inches(0.6),
   "消融实验: Noise Mode的关键作用", 36, WHITE, True)
card(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.0),
   "With Noise Mode", "Cross Acc = 89.71%", "Neural Cleanse检测不到", GREEN)
card(s, Inches(6.8), Inches(2.0), Inches(5.5), Inches(2.0),
   "Without Noise Mode", "Cross Acc = 0.00%", "Neural Cleanse能检测到!", RED)
tb(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.0),
   "结论: Noise Mode是WaNet绕过防御的决定性因素\n"
   "没有Noise Mode -> 模型学到像素级捷径 -> 被Neural Cleanse发现\n"
   "有Noise Mode -> 模型被迫学习真正的warping变换 -> 防御失效",
   17, WHITE, a=PP_ALIGN.CENTER)

# === SLIDE 7: s & k Ablation ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(0.8), Inches(0.8), Inches(2.5))
tb(s, Inches(0.8), Inches(1.0), Inches(7), Inches(0.6),
   "消融实验: Warping参数敏感性分析", 36, WHITE, True)
tb(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.4),
   "Warping强度 s 的影响", 22, GOLD, True)
card(s, Inches(0.8), Inches(2.4), Inches(2.6), Inches(1.5), "s = 0.25", "ASR 74%", "太弱->攻击失效", RED)
card(s, Inches(3.7), Inches(2.4), Inches(2.6), Inches(1.5), "s = 0.50", "ASR 99.59%", "最优参数", GREEN)
card(s, Inches(6.6), Inches(2.4), Inches(2.6), Inches(1.5), "s = 1.00", "ASR 98.29%", "性能饱和", GRAY)
tb(s, Inches(0.8), Inches(4.2), Inches(5.5), Inches(0.4),
   "网格大小 k 的影响", 22, GOLD, True)
card(s, Inches(0.8), Inches(4.8), Inches(2.6), Inches(1.5), "k = 2", "ASR 92.82%", "太粗->被当噪声", RED)
card(s, Inches(3.7), Inches(4.8), Inches(2.6), Inches(1.5), "k = 4", "ASR 99.15%", "最优参数", GREEN)
card(s, Inches(6.6), Inches(4.8), Inches(2.6), Inches(1.5), "k = 8", "ASR 97.55%", "性能饱和", GRAY)
tb(s, Inches(10), Inches(2.0), Inches(2.8), Inches(4.0),
   "关键发现\n\nWaNet对参数\n高度敏感\n\ns<0.5或k<4时\n攻击大幅下降\n\n存在狭窄的\n有效参数区间",
   14, LGOLD, a=PP_ALIGN.CENTER)

# === SLIDE 8: Insight ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(0.8), Inches(0.8), Inches(2.5))
tb(s, Inches(0.8), Inches(1.0), Inches(7), Inches(0.6),
   "核心洞察: Warping在频域留下痕迹", 36, WHITE, True)
mtb(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.0), [
    ("消融实验揭示了一个重要线索:", False, GRAY),
    ("Warping = grid_sample + bilinear interpolation", True, GOLD),
    ("", False, GRAY),
    ("插值运算 = 信号处理 -> 必然在频域留下痕迹!", True, WHITE),
    ("", False, GRAY),
    ("这个发现指向了一个全新的检测思路:", False, GRAY),
    ("", False, GRAY),
    ("现有防御: 看模型输出(logits/梯度) -> 检测不到WaNet", False, RED),
    ("我们的方法: 看图像频域(FFT频谱) -> 能检测到!", True, GREEN),
], 18)
tb(s, Inches(7.5), Inches(2.0), Inches(5), Inches(3.0),
   "Warping的本质\n\ngrid_sample对图像\n做双线性插值\n\n每个插值点 = 周围像素\n的加权平均\n\n-> 频域高频分量异常",
   16, LGOLD, a=PP_ALIGN.CENTER)
line(s, Inches(0.8), Inches(5.5), Inches(11.7))
tb(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.8),
   "这是一个论文没有探讨过的维度: 不可感知 != 不可检测", 18, GOLD, True, PP_ALIGN.CENTER)

# === SLIDE 9: Detector Design ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(0.8), Inches(0.8), Inches(2.5))
tb(s, Inches(0.8), Inches(1.0), Inches(7), Inches(0.6),
   "频域Warping检测器设计", 36, WHITE, True)
steps = [
    ("Step 1: 生成图像对", "Clean + Warped\n各500张", GRAY),
    ("Step 2: FFT变换", "傅里叶变换\n->频谱图", GRAY),
    ("Step 3: 特征提取", "高频能量比\n频谱熵等6维", GRAY),
    ("Step 4: SVM分类", "RBF核SVM\n30%测试集", GRAY),
    ("Step 5: 评估", "ROC曲线\nvs STRIP", GOLD),
]
for i, (t, d, ac) in enumerate(steps):
    card(s, Inches(0.8+i*2.5), Inches(2.0), Inches(2.2), Inches(2.0), t, d, "", ac)
    if i < 4:
        tb(s, Inches(0.8+i*2.5+2.3), Inches(2.8), Inches(0.3), Inches(0.3),
           "->", 22, GOLD, a=PP_ALIGN.CENTER)
card(s, Inches(0.8), Inches(4.5), Inches(5.5), Inches(1.5),
   "6维频域特征", "高频能量比 / 频谱熵 / 均值频率 / 方差频率 / 低频能量比 / 中频能量比", "", TEAL)
card(s, Inches(6.8), Inches(4.5), Inches(5.5), Inches(1.5),
   "检测结果", "准确率 67.33%  |  AUC 0.735  |  STRIP AUC=0.50", "", GREEN)

# === SLIDE 10: FFT comparison ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(0.8), Inches(0.8), Inches(2.5))
tb(s, Inches(0.8), Inches(1.0), Inches(7), Inches(0.6),
   "频域可视化: Clean vs Warped 频谱对比", 36, WHITE, True)
for i, label in enumerate([
    "Clean Image", "Clean FFT Spectrum", "Clean Radial Profile",
    "Warped Image (Backdoor)", "Warped FFT Spectrum", "FFT Difference Map"
]):
    x = Inches(0.8 + (i%3)*4.0); y = Inches(1.8 + (i//3)*2.6)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(2.3))
    sh.fill.solid(); sh.fill.fore_color.rgb = DARK
    sh.line.color.rgb = GOLD if i==5 else GRAY; sh.line.width = Pt(1)
    tb(s, x+Pt(8), y+Pt(4), Inches(3.5), Inches(0.3), label, 13,
       GOLD if i==5 else GRAY, a=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4),
   "插入 results/freq_detector/fft_spectrum_comparison.png  |  注意右下角: 差异图有明显结构化异常!",
   14, GOLD, a=PP_ALIGN.CENTER)

# === SLIDE 11: ROC ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(0.8), Inches(0.8), Inches(2.5))
tb(s, Inches(0.8), Inches(1.0), Inches(7), Inches(0.6),
   "检测性能对比: 频域检测器 vs STRIP", 36, WHITE, True)
card(s, Inches(0.8), Inches(2.0), Inches(3.8), Inches(2.0),
   "频域检测器 (Ours)", "AUC = 0.735", "准确率 67.33%", GREEN)
card(s, Inches(5.0), Inches(2.0), Inches(3.8), Inches(2.0),
   "STRIP", "AUC = 0.50", "等同随机猜测", RED)
card(s, Inches(9.2), Inches(2.0), Inches(3.5), Inches(2.0),
   "提升", "+23.5% AUC", "远超随机基线", GOLD)
tb(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.4),
   "插入 results/freq_detector/roc_curve.png", 14, GRAY, a=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.6),
   "STRIP看模型输出熵 -> 无法区分Clean/Warped (AUC=0.50)\n我们的频域检测器看图像信号 -> 能有效检测 (AUC=0.735)",
   16, WHITE, a=PP_ALIGN.CENTER)

# === SLIDE 12: Full Results Table ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(0.8), Inches(0.8), Inches(2.5))
tb(s, Inches(0.8), Inches(1.0), Inches(7), Inches(0.6),
   "全部实验结果汇总", 36, WHITE, True)
data = [
    ["Epochs", "Experiment", "Clean Acc", "ASR", "Cross Acc", "Key Finding"],
    ["400", "Noise WITH (baseline)", "92.84", "99.15", "89.71", "Baseline"],
    ["400", "Noise WITHOUT", "92.74", "98.21", "0.00", "Noise is KEY!"],
    ["400", "s=0.25", "92.39", "74.12", "89.38", "s too small -> fail"],
    ["400", "s=0.50 (best)", "92.81", "99.59", "89.69", "Optimal"],
    ["400", "s=1.00", "92.82", "98.29", "89.30", "Saturation"],
    ["80", "k=2", "92.14", "92.82", "88.95", "k too small -> weak"],
    ["400", "k=4 (best)", "92.84", "99.15", "89.71", "Optimal"],
    ["80", "k=8", "91.94", "97.55", "88.04", "Saturation"],
    ["400", "Seed=0", "92.81", "99.59", "89.69", "Stable"],
    ["80", "Seed=1", "92.10", "98.70", "88.72", "Stable"],
]
r, c = len(data), len(data[0])
tbl = s.shapes.add_table(r, c, Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.8)).table
for ri in range(r):
    for ci in range(c):
        cell = tbl.cell(ri, ci); cell.text = data[ri][ci]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11); p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
            if ri == 0: p.font.bold = True; p.font.color.rgb = NAVY
            elif ri in [1, 6]: p.font.color.rgb = WHITE
            else: p.font.color.rgb = GRAY
        if ri == 0: cell.fill.solid(); cell.fill.fore_color.rgb = GOLD
        elif ri in [1, 6]: cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1A, 0x2D, 0x50)

# === SLIDE 13: Conclusion ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(2), Inches(1.5), Inches(9.333))
tb(s, Inches(2), Inches(1.8), Inches(9.333), Inches(1.0),
   "结论", 44, WHITE, True, PP_ALIGN.CENTER)
conclusions = [
    ("1", "WaNet攻击高度有效", "CIFAR-10上ASR达99.15%, 同时保持92.84%正常精度", GREEN),
    ("2", "Noise Mode是绕过防御的关键", "无Noise时Cross Acc=0%, 后门可被Neural Cleanse检测", GOLD),
    ("3", "Warping参数高度敏感", "s<0.5或k<4时攻击效果大幅下降 -- 存在狭窄有效区间", RED),
    ("4", "频域检测: 我们的创新", "首次从图像信号层面检测warping后门, AUC 0.735远超STRIP 0.50", GREEN),
]
for i, (num, title, desc, color) in enumerate(conclusions):
    y = Inches(3.0 + i*1.0)
    tb(s, Inches(2.0), y, Inches(0.5), Inches(0.5), num, 28, color, True)
    tb(s, Inches(2.7), y, Inches(3.5), Inches(0.4), title, 22, WHITE, True)
    tb(s, Inches(2.7), y+Inches(0.45), Inches(8), Inches(0.4), desc, 15, GRAY)
line(s, Inches(2), Inches(6.8), Inches(9.333))

# === SLIDE 14: Thank you ===
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(2), Inches(3.0), Inches(9.333))
tb(s, Inches(2), Inches(3.3), Inches(9.333), Inches(1.0),
   "感谢聆听", 52, WHITE, True, PP_ALIGN.CENTER)
tb(s, Inches(2), Inches(4.5), Inches(9.333), Inches(0.6),
   "WaNet复现与分析 | AI安全课程设计", 18, GOLD, a=PP_ALIGN.CENTER)
tb(s, Inches(2), Inches(5.3), Inches(9.333), Inches(0.5),
   "github.com/axgo-ja/NIS3353", 16, GRAY, a=PP_ALIGN.CENTER)
line(s, Inches(2), Inches(6.2), Inches(9.333))

prs.save(OUT)
print(f'PPT saved: {OUT}')
print(f'Slides: {len(prs.slides)}')
